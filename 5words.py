import os
from collections import Counter
import spacy
from spacy.lang.en.stop_words import STOP_WORDS
from pptx import Presentation
from docx import Document
import textract

# Load English language model
nlp = spacy.load("en_core_web_sm")

def extract_text(file_path):
    _, file_extension = os.path.splitext(file_path)
    if file_extension.lower() == '.pdf':
        print("PDF files are not supported.")
        return None

    text = ""
    if file_extension.lower() == '.docx':
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    elif file_extension.lower() == '.doc':
        try:
            text = textract.process(file_path).decode('utf-8')
        except Exception as e:
            print(f"Failed to extract text: {e}")
            return None
    elif file_extension.lower() in ['.pptx', '.ppt']:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    else:
        print(f"Unsupported file format: {file_extension}")
        return None

    return text

def generate_representative_word(text):
    doc = nlp(text)

    # Extract nouns and lemmatize
    nouns = [token.lemma_ for token in doc if token.pos_ == "NOUN" and token.text.lower() not in STOP_WORDS]

    # Count the frequency of nouns
    noun_freq = Counter(nouns)

    # Find the most common nouns
    most_common_nouns = noun_freq.most_common()

    # Take top 5 most frequent nouns as higher-level concepts
    higher_level_concepts = [noun for noun, _ in most_common_nouns[:5]]

    return higher_level_concepts

def process_file(file_path, file_keywords):
    print(f"Processing file: {file_path}")
    text = extract_text(file_path)
    if text:
        representative_words = generate_representative_word(text)
        if representative_words:
            file_keywords[file_path] = representative_words
        else:
            print("No representative words found for the file.")
        print("----------------------------------------------")

def process_directory(directory_path, file_keywords):
    if os.path.isfile(directory_path):
        process_file(directory_path, file_keywords)
    elif os.path.isdir(directory_path):
        for root, _, files in os.walk(directory_path):
            for filename in files:
                file_path = os.path.join(root, filename)
                process_file(file_path, file_keywords)
    else:
        print(f"Invalid path: {directory_path}")

if _name_ == "_main_":
    path = input("Enter the file or directory path: ")
    file_keywords = {}
    process_directory(path, file_keywords)

    # Display file paths and their corresponding representative keywords
    print("File Paths and Their Representative Keywords:")
    for file_path, keywords in file_keywords.items():
        print(f"File: {file_path} - Keywords: {keywords}")
